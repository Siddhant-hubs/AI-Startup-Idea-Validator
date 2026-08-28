from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

slides = [
    ("AI Startup Idea Validator", ["Milestone 4 — Final Demonstration", "Market → SWOT/Risk → MVP → GTM → Advisor → PDF"]),
    ("Problem Statement", ["Founders need fast evidence before building.", "Manual market, competitor, product and launch analysis is fragmented.", "The system connects these decisions in one workflow."]),
    ("System Architecture", ["FastAPI frontend/API", "CrewAI market research with Tavily + competitor retrieval", "SWOT/Risk → MVP/MoSCoW → GTM", "Knowledge Base → Conversational Advisor", "Report Generation → downloadable PDF"]),
    ("Market Validation", ["Live web research", "Market size and growth evidence", "Real competitor discovery", "Confidence score and verified source URLs"]),
    ("SWOT & Risk Agent", ["Strengths / Weaknesses / Opportunities / Threats", "High, medium and low risks", "Competitor risk and market-demand prediction", "Mitigation recommendations"]),
    ("MVP Recommendation", ["MoSCoW prioritization", "Must Have / Should Have / Could Have / Won't Have", "Effort and impact", "Tech stack, timeline, resources and success metrics"]),
    ("Go-To-Market Strategy", ["Positioning and value proposition", "Primary and secondary channels", "First 100 and first 1000 customers", "Pricing, launch plan, partnerships and growth metrics"]),
    ("Conversational Advisor", ["Validation outputs are ingested into a session knowledge base.", "Follow-up questions reference SWOT, MVP and GTM.", "Conversation history is passed to the advisor."]),
    ("Milestone 4: Report + Testing", ["Structured downloadable PDF", "Output-contract quality checks", "Mocked end-to-end API tests", "Search-query optimization and prompt-engineering documentation"]),
    ("Live Demo Flow", ["Enter startup idea", "Run Full Validation", "Review five tabs", "Download PDF", "Ask advisor follow-up questions", "Run python run_tests.py"]),
    ("Limitations & Future Scope", ["LLM strategy is decision support, not a guarantee.", "Market claims should be checked against sources.", "Current chatbot KB is session-scoped.", "Future: persistence, authentication, source verification and production observability."]),
    ("Thank You / Q&A", ["Demo repository + deployed URL are environment-specific.", "Questions?"])
]

def add_slide(title, bullets, first=False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid(); bg.fore_color.rgb = __import__('pptx').dml.color.RGBColor(248,247,252)
    # title
    box = slide.shapes.add_textbox(Inches(.7), Inches(.55), Inches(12), Inches(1))
    tf = box.text_frame
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(30); p.font.bold = True
    p.font.color.rgb = __import__('pptx').dml.color.RGBColor(91,54,190)
    # accent
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(.7), Inches(1.55), Inches(1.2), Inches(.08))
    line.fill.solid(); line.fill.fore_color.rgb = __import__('pptx').dml.color.RGBColor(91,54,190); line.line.fill.background()
    # bullets
    box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11), Inches(4.8))
    tf = box.text_frame; tf.word_wrap = True
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b; p.font.size = Pt(20 if not first else 22); p.space_after = Pt(15)
    return slide

for i,(title,bullets) in enumerate(slides):
    add_slide(title, bullets, first=(i==0))

out='/mnt/data/AI-Startup-Idea-Validator-Milestone3/AI_Startup_Validator_Milestone4_Demo.pptx'
prs.save(out)
print(out)
